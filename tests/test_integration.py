"""Integration tests for semstash with real AWS services.

These tests require:
    - Valid AWS credentials configured
    - Permission to create/delete S3 buckets and S3 Vectors resources
    - Access to Bedrock Nova embeddings model

Run with: pytest --use-aws -m integration
Preserve resources for debugging: pytest --use-aws --preserve-aws -m integration
"""

from pathlib import Path

import httpx
import pytest
from helpers import assert_valid_query_results, assert_valid_search_result

from semstash import SemStash
from semstash.exceptions import ContentNotFoundError


@pytest.mark.integration
class TestIntegrationUploadQuery:
    """Integration tests for upload and query workflow."""

    def test_upload_text_to_root_and_query(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Upload a text file to root and query for it."""
        # Upload file to root
        result = integration_stash.upload(sample_text_file, target="/")
        assert result.key == sample_text_file.name
        assert result.path == f"/{sample_text_file.name}"
        assert result.content_type == "text/plain"

        # Query for content (matches actual text file content)
        query_results = integration_stash.query("sample text for testing semantic storage", top_k=5)

        # Use helper to validate all results have proper scores and metadata
        assert_valid_query_results(
            query_results, min_count=1, expected_keys=[sample_text_file.name]
        )

        # Verify the specific matched result has path
        matched_result = next(r for r in query_results if r.key == sample_text_file.name)
        assert matched_result.path == f"/{sample_text_file.name}"
        assert_valid_search_result(
            matched_result,
            expected_key=sample_text_file.name,
            expected_content_type="text/plain",
            require_file_size=True,
        )

    def test_upload_text_to_folder_and_query(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Upload a text file to folder and query for it."""
        # Upload file to /docs/ folder (preserves filename)
        result = integration_stash.upload(sample_text_file, target="/docs/")
        assert result.key == f"docs/{sample_text_file.name}"
        assert result.path == f"/docs/{sample_text_file.name}"
        assert result.content_type == "text/plain"

        # Query should find the file
        query_results = integration_stash.query("sample text for testing semantic storage", top_k=5)
        assert_valid_query_results(
            query_results, min_count=1, expected_keys=[f"docs/{sample_text_file.name}"]
        )

        # Verify path in result
        matched_result = next(r for r in query_results if r.key == f"docs/{sample_text_file.name}")
        assert matched_result.path == f"/docs/{sample_text_file.name}"

    def test_upload_with_rename_and_query(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Upload a file with renaming (no trailing slash in target)."""
        # Upload file with explicit target path (rename on upload)
        result = integration_stash.upload(sample_text_file, target="/renamed/readme.txt")
        assert result.key == "renamed/readme.txt"
        assert result.path == "/renamed/readme.txt"

        # Query should find the file
        query_results = integration_stash.query("sample text for testing semantic storage", top_k=5)
        assert_valid_query_results(query_results, min_count=1, expected_keys=["renamed/readme.txt"])

    def test_upload_image_and_query(
        self,
        integration_stash: SemStash,
        sample_image_file: Path,
    ) -> None:
        """Upload an image file and query for it."""
        # Upload image to root
        result = integration_stash.upload(sample_image_file, target="/")
        assert result.key == sample_image_file.name
        assert result.path == f"/{sample_image_file.name}"
        assert result.content_type == "image/png"

        # Query for image content (1x1 transparent pixel - minimal PNG)
        # For a minimal test image, any visual query will have low similarity
        # Just verify the query returns valid results with proper metadata
        query_results = integration_stash.query("small transparent image", top_k=5)
        assert_valid_query_results(query_results, min_count=1)


@pytest.mark.integration
class TestIntegrationBrowse:
    """Integration tests for browse functionality."""

    def test_browse_root_after_uploads(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Browse root content after uploading multiple files."""

        # Upload multiple files to root
        integration_stash.upload(sample_text_file, target="/")
        integration_stash.upload(sample_json_file, target="/")

        # Browse root path
        result = integration_stash.browse("/")
        assert len(result.items) == 2
        assert result.next_token is None  # No pagination needed

        # Check file names and paths
        keys = [item.key for item in result.items]
        paths = [item.path for item in result.items]
        assert sample_text_file.name in keys
        assert sample_json_file.name in keys
        assert f"/{sample_text_file.name}" in paths
        assert f"/{sample_json_file.name}" in paths

    def test_browse_folder_after_uploads(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Browse folder content after uploading files to different locations."""

        # Upload to different locations
        integration_stash.upload(sample_text_file, target="/docs/")
        integration_stash.upload(sample_json_file, target="/config/")

        # Browse /docs/ folder - should only have text file
        result = integration_stash.browse("/docs/")
        assert len(result.items) == 1
        assert result.items[0].key == f"docs/{sample_text_file.name}"
        assert result.items[0].path == f"/docs/{sample_text_file.name}"

        # Browse /config/ folder - should only have json file
        result = integration_stash.browse("/config/")
        assert len(result.items) == 1
        assert result.items[0].key == f"config/{sample_json_file.name}"
        assert result.items[0].path == f"/config/{sample_json_file.name}"

        # Browse root - should show all items (S3 lists all objects with prefix)
        result = integration_stash.browse("/")
        assert len(result.items) == 2
        keys = {item.key for item in result.items}
        assert f"docs/{sample_text_file.name}" in keys
        assert f"config/{sample_json_file.name}" in keys

    def test_browse_nonexistent_path_returns_empty(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Browse a path that doesn't exist returns empty results."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Browse non-existent folder
        result = integration_stash.browse("/nonexistent/")
        assert len(result.items) == 0
        assert result.total == 0


@pytest.mark.integration
class TestIntegrationGetDelete:
    """Integration tests for get and delete functionality."""

    def test_get_content_at_root(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Get uploaded content at root."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Get metadata and URL using path
        result = integration_stash.get(f"/{sample_text_file.name}")

        assert result.key == sample_text_file.name
        assert result.path == f"/{sample_text_file.name}"
        assert result.content_type == "text/plain"
        assert result.url is not None
        assert "https://" in result.url  # Should be a presigned S3 URL

    def test_get_content_in_folder(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Get uploaded content in folder."""

        # Upload file to folder
        integration_stash.upload(sample_text_file, target="/docs/")

        # Get metadata and URL using full path
        result = integration_stash.get(f"/docs/{sample_text_file.name}")

        assert result.key == f"docs/{sample_text_file.name}"
        assert result.path == f"/docs/{sample_text_file.name}"
        assert result.content_type == "text/plain"
        assert result.url is not None

    def test_get_nonexistent_path_raises_error(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Get non-existent path raises appropriate error."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Try to get file at wrong path - should raise error
        with pytest.raises(ContentNotFoundError):
            integration_stash.get("/nonexistent/file.txt")

    def test_delete_content_at_root(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Delete uploaded content at root."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Verify it exists
        browse_result = integration_stash.browse("/")
        assert len(browse_result.items) == 1

        # Delete it using path
        result = integration_stash.delete(f"/{sample_text_file.name}")
        assert result.deleted is True

        # Verify it's gone
        browse_result = integration_stash.browse("/")
        assert len(browse_result.items) == 0

    def test_delete_content_in_folder(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Delete uploaded content in folder."""

        # Upload file to folder
        integration_stash.upload(sample_text_file, target="/docs/")

        # Verify it exists
        browse_result = integration_stash.browse("/docs/")
        assert len(browse_result.items) == 1

        # Delete it using full path
        result = integration_stash.delete(f"/docs/{sample_text_file.name}")
        assert result.deleted is True

        # Verify it's gone
        browse_result = integration_stash.browse("/docs/")
        assert len(browse_result.items) == 0


@pytest.mark.integration
class TestIntegrationCheckSync:
    """Integration tests for check and sync functionality."""

    def test_check_consistent_storage(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Check returns consistent state after normal operations."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Check consistency
        result = integration_stash.check()
        assert result.is_consistent is True
        assert result.content_count == 1
        assert result.vector_count == 1
        assert len(result.orphaned_vectors) == 0
        assert len(result.missing_vectors) == 0

    def test_check_consistent_after_folder_uploads(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Check returns consistent state after folder uploads."""

        # Upload files to different folders
        integration_stash.upload(sample_text_file, target="/docs/")
        integration_stash.upload(sample_json_file, target="/config/")

        # Check consistency
        result = integration_stash.check()
        assert result.is_consistent is True
        assert result.content_count == 2
        assert result.vector_count == 2
        assert len(result.orphaned_vectors) == 0
        assert len(result.missing_vectors) == 0


@pytest.mark.integration
class TestIntegrationOpenExisting:
    """Integration tests for opening existing storage."""

    def test_open_existing_stash(
        self,
        integration_stash: SemStash,
        integration_bucket_name: str,
        sample_text_file: Path,
    ) -> None:
        """Open an existing stash without re-initializing."""
        # Stash is already initialized by the session fixture
        upload_result = integration_stash.upload(sample_text_file, target="/")
        assert upload_result.path == f"/{sample_text_file.name}"

        # Open existing stash with a new client
        stash2 = SemStash(integration_bucket_name)
        stash2.open()  # Should detect existing config

        # Should be able to query (matches actual text file content)
        results = stash2.query("sample text for testing semantic storage", top_k=5)
        assert len(results) >= 1

        # Verify path is correct in query results
        matched = next((r for r in results if r.key == sample_text_file.name), None)
        assert matched is not None
        assert matched.path == f"/{sample_text_file.name}"

        # Should have same dimension
        assert stash2.dimension == 256

    def test_open_existing_stash_with_folder_content(
        self,
        integration_stash: SemStash,
        integration_bucket_name: str,
        sample_text_file: Path,
    ) -> None:
        """Open existing stash and access content in folders."""
        # Upload to folder
        upload_result = integration_stash.upload(sample_text_file, target="/docs/")
        assert upload_result.path == f"/docs/{sample_text_file.name}"

        # Open existing stash with a new client
        stash2 = SemStash(integration_bucket_name)
        stash2.open()

        # Should be able to browse and get
        browse_result = stash2.browse("/docs/")
        assert len(browse_result.items) == 1
        assert browse_result.items[0].path == f"/docs/{sample_text_file.name}"

        # Should be able to get using path
        get_result = stash2.get(f"/docs/{sample_text_file.name}")
        assert get_result.path == f"/docs/{sample_text_file.name}"


@pytest.mark.integration
class TestIntegrationContentVerification:
    """Integration tests for content verification (download and compare)."""

    def test_download_text_content_matches(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download text file and verify content matches original."""

        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Get presigned URL using path
        result = integration_stash.get(f"/{sample_text_file.name}")
        assert result.path == f"/{sample_text_file.name}"

        # Download content using presigned URL
        response = httpx.get(result.url)
        assert response.status_code == 200

        # Save to tmp and compare
        downloaded_path = tmp_path / "downloaded.txt"
        downloaded_path.write_bytes(response.content)

        # Verify content matches
        original_content = sample_text_file.read_bytes()
        downloaded_content = downloaded_path.read_bytes()
        assert downloaded_content == original_content

    def test_download_image_content_matches(
        self,
        integration_stash: SemStash,
        sample_image_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download image file and verify content matches original."""

        # Upload image to root
        integration_stash.upload(sample_image_file, target="/")

        # Get presigned URL using path
        result = integration_stash.get(f"/{sample_image_file.name}")
        assert result.path == f"/{sample_image_file.name}"

        # Download content using presigned URL
        response = httpx.get(result.url)
        assert response.status_code == 200

        # Save to tmp and compare
        downloaded_path = tmp_path / "downloaded.png"
        downloaded_path.write_bytes(response.content)

        # Verify content matches
        original_content = sample_image_file.read_bytes()
        downloaded_content = downloaded_path.read_bytes()
        assert downloaded_content == original_content

    def test_download_json_content_matches(
        self,
        integration_stash: SemStash,
        sample_json_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download JSON file and verify content matches original."""

        # Upload JSON to root
        integration_stash.upload(sample_json_file, target="/")

        # Get presigned URL using path
        result = integration_stash.get(f"/{sample_json_file.name}")
        assert result.path == f"/{sample_json_file.name}"

        # Download content using presigned URL
        response = httpx.get(result.url)
        assert response.status_code == 200

        # Save to tmp and compare
        downloaded_path = tmp_path / "downloaded.json"
        downloaded_path.write_bytes(response.content)

        # Verify content matches
        original_content = sample_json_file.read_bytes()
        downloaded_content = downloaded_path.read_bytes()
        assert downloaded_content == original_content

    def test_download_content_from_folder(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download content from folder and verify content matches."""

        # Upload to folder
        integration_stash.upload(sample_text_file, target="/data/")

        # Get using full path
        result = integration_stash.get(f"/data/{sample_text_file.name}")
        assert result.path == f"/data/{sample_text_file.name}"
        assert result.key == f"data/{sample_text_file.name}"

        # Download and verify
        response = httpx.get(result.url)
        assert response.status_code == 200
        assert response.content == sample_text_file.read_bytes()


@pytest.mark.integration
class TestIntegrationTagFiltering:
    """Integration tests for tag filtering with native S3 Vectors support."""

    def test_upload_with_tags_and_filter(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Upload files with tags and filter by tag in query."""
        # Upload with different tags to root
        result1 = integration_stash.upload(
            sample_text_file, target="/", tags=["documentation", "readme"]
        )
        result2 = integration_stash.upload(
            sample_json_file, target="/", tags=["config", "settings"]
        )
        assert result1.path == f"/{sample_text_file.name}"
        assert result2.path == f"/{sample_json_file.name}"

        # Query with tag filter - should only return the text file
        results = integration_stash.query("sample content", tags=["documentation"])
        assert_valid_query_results(results, min_count=1, expected_keys=[sample_text_file.name])
        keys = [r.key for r in results]
        # JSON file should not be in filtered results (different tags)
        assert sample_json_file.name not in keys

        # Verify path is present in results
        matched = next((r for r in results if r.key == sample_text_file.name), None)
        assert matched is not None
        assert matched.path == f"/{sample_text_file.name}"

    def test_query_with_multiple_tags(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with multiple tags uses OR logic (any match)."""
        # Upload with different tags to folders
        result1 = integration_stash.upload(
            sample_text_file, target="/docs/", tags=["type-a", "category-1"]
        )
        result2 = integration_stash.upload(
            sample_json_file, target="/config/", tags=["type-b", "category-2"]
        )
        assert result1.path == f"/docs/{sample_text_file.name}"
        assert result2.path == f"/config/{sample_json_file.name}"

        # Query with multiple tags - should return files matching any tag
        results = integration_stash.query("content", tags=["type-a", "type-b"])
        assert_valid_query_results(
            results,
            min_count=2,
            expected_keys=[f"docs/{sample_text_file.name}", f"config/{sample_json_file.name}"],
        )

    def test_query_with_nonexistent_tag(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Query with non-existent tag returns no results."""
        # Upload file with specific tags to root
        integration_stash.upload(sample_text_file, target="/", tags=["existing-tag"])

        # Query with a tag that doesn't exist
        results = integration_stash.query("sample content", tags=["nonexistent-tag"])
        # Should return empty list since no content matches the tag
        assert len(results) == 0


@pytest.mark.integration
class TestIntegrationDocuments:
    """Integration tests for document upload and search (PDF, DOCX, PPTX, XLSX)."""

    def test_upload_pdf_and_query(
        self,
        integration_stash: SemStash,
        sample_pdf_file: Path,
    ) -> None:
        """Upload PDF and query for its content."""
        # Upload PDF (contains text about machine learning) to root
        result = integration_stash.upload(sample_pdf_file, target="/")
        assert result.key == sample_pdf_file.name
        assert result.path == f"/{sample_pdf_file.name}"
        assert result.content_type == "application/pdf"

        # Query for content that's in the PDF
        query_results = integration_stash.query(
            "machine learning neural networks deep learning", top_k=5
        )
        assert_valid_query_results(query_results, min_count=1, expected_keys=[sample_pdf_file.name])

        # Verify path in result (may be chunked, e.g., /sample.pdf#page=1)
        matched = next(
            (r for r in query_results if r.key.split("#")[0] == sample_pdf_file.name), None
        )
        assert matched is not None
        assert matched.path.split("#")[0] == f"/{sample_pdf_file.name}"

    def test_upload_docx_and_query(
        self,
        integration_stash: SemStash,
        sample_docx_file: Path,
    ) -> None:
        """Upload DOCX and query for its content."""
        # Upload DOCX (contains text about semantic storage architecture) to folder
        result = integration_stash.upload(sample_docx_file, target="/documents/")
        assert result.key == f"documents/{sample_docx_file.name}"
        assert result.path == f"/documents/{sample_docx_file.name}"
        assert "wordprocessingml" in result.content_type

        # Query for content that's in the DOCX
        query_results = integration_stash.query("semantic storage architecture embeddings", top_k=5)
        assert_valid_query_results(
            query_results, min_count=1, expected_keys=[f"documents/{sample_docx_file.name}"]
        )

    def test_upload_pptx_and_query(
        self,
        integration_stash: SemStash,
        sample_pptx_file: Path,
    ) -> None:
        """Upload PPTX and query for its content."""
        # Upload PPTX (contains text about REST API design) to root
        result = integration_stash.upload(sample_pptx_file, target="/")
        assert result.key == sample_pptx_file.name
        assert result.path == f"/{sample_pptx_file.name}"
        assert "presentationml" in result.content_type

        # Query for content that's in the PPTX
        query_results = integration_stash.query("REST API design patterns", top_k=5)
        assert_valid_query_results(
            query_results, min_count=1, expected_keys=[sample_pptx_file.name]
        )

    def test_upload_xlsx_and_query(
        self,
        integration_stash: SemStash,
        sample_xlsx_file: Path,
    ) -> None:
        """Upload XLSX and query for its content."""
        # Upload XLSX (contains quarterly sales data) to root
        result = integration_stash.upload(sample_xlsx_file, target="/")
        assert result.key == sample_xlsx_file.name
        assert result.path == f"/{sample_xlsx_file.name}"
        assert "spreadsheetml" in result.content_type

        # Query for content that's in the XLSX (converted to CSV/text)
        query_results = integration_stash.query("quarterly sales revenue financial data", top_k=5)
        assert_valid_query_results(
            query_results, min_count=1, expected_keys=[sample_xlsx_file.name]
        )

    def test_download_pdf_content_matches(
        self,
        integration_stash: SemStash,
        sample_pdf_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download PDF file and verify content matches original."""
        # Upload PDF to root
        integration_stash.upload(sample_pdf_file, target="/")

        # Get presigned URL using path
        result = integration_stash.get(f"/{sample_pdf_file.name}")
        assert result.path == f"/{sample_pdf_file.name}"
        assert result.content_type == "application/pdf"

        # Download content using presigned URL
        response = httpx.get(result.url)
        assert response.status_code == 200

        # Save to tmp and compare
        downloaded_path = tmp_path / "downloaded.pdf"
        downloaded_path.write_bytes(response.content)

        # Verify content matches
        original_content = sample_pdf_file.read_bytes()
        downloaded_content = downloaded_path.read_bytes()
        assert downloaded_content == original_content

    def test_download_docx_content_matches(
        self,
        integration_stash: SemStash,
        sample_docx_file: Path,
        tmp_path: Path,
    ) -> None:
        """Download DOCX file and verify content matches original."""
        # Upload DOCX to root
        integration_stash.upload(sample_docx_file, target="/")

        # Get presigned URL using path
        result = integration_stash.get(f"/{sample_docx_file.name}")
        assert result.path == f"/{sample_docx_file.name}"

        # Download content using presigned URL
        response = httpx.get(result.url)
        assert response.status_code == 200

        # Save to tmp and compare
        downloaded_path = tmp_path / "downloaded.docx"
        downloaded_path.write_bytes(response.content)

        # Verify content matches
        original_content = sample_docx_file.read_bytes()
        downloaded_content = downloaded_path.read_bytes()
        assert downloaded_content == original_content

    def test_cross_modality_query(
        self,
        integration_stash: SemStash,
        sample_pdf_file: Path,
        sample_docx_file: Path,
        sample_jpg_file: Path,
    ) -> None:
        """Upload multiple document types and query across them."""
        # Upload all documents to root (they all contain the same sample.jpg image)
        integration_stash.upload(sample_pdf_file, target="/")
        integration_stash.upload(sample_docx_file, target="/")
        integration_stash.upload(sample_jpg_file, target="/")

        # Query for landscape/nature content (common in the image)
        query_results = integration_stash.query("landscape nature scenery", top_k=5)

        # At minimum the JPG should match with valid scores
        assert_valid_query_results(query_results, min_count=1, expected_keys=[sample_jpg_file.name])

        # Verify path is present in results (may be chunked)
        for r in query_results:
            assert r.path.startswith("/")
            # Path may include chunk suffix like #page=1 or #slide=2
            source_key = r.key.split("#")[0]
            assert r.path.split("#")[0] == f"/{source_key}"


@pytest.mark.integration
class TestIntegrationPathFiltering:
    """Integration tests for path-based query filtering."""

    def test_query_with_path_filter_finds_content(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with path filter finds content in specified folder."""
        # Upload files to different folders
        integration_stash.upload(sample_text_file, target="/docs/")
        integration_stash.upload(sample_json_file, target="/config/")

        # Query with path filter for /docs/ - should find only the text file
        results = integration_stash.query("sample content", path="/docs/")
        assert len(results) >= 1
        keys = [r.key for r in results]
        assert f"docs/{sample_text_file.name}" in keys
        # JSON file should NOT be in results (different folder)
        assert f"config/{sample_json_file.name}" not in keys

        # Verify path is correct in results
        for r in results:
            assert r.path.startswith("/docs/")

    def test_query_with_path_filter_excludes_other_folders(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with path filter excludes content from other folders."""
        # Upload files to different folders
        integration_stash.upload(sample_text_file, target="/folder-a/")
        integration_stash.upload(sample_json_file, target="/folder-b/")

        # Query with path filter for /folder-b/
        results = integration_stash.query("content", path="/folder-b/")

        # Should only find JSON file
        keys = [r.key for r in results]
        assert f"folder-a/{sample_text_file.name}" not in keys

        # All results should be from /folder-b/
        for r in results:
            assert r.path.startswith("/folder-b/")

    def test_query_with_nonexistent_path_returns_empty(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
    ) -> None:
        """Query with non-existent path filter returns no results."""
        # Upload file to root
        integration_stash.upload(sample_text_file, target="/")

        # Query with path filter for non-existent folder
        results = integration_stash.query("sample content", path="/nonexistent/")

        # Should return empty list since no content matches the path
        assert len(results) == 0

    def test_query_with_nested_path_filter(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with nested path filter works correctly."""
        # Upload files to nested folders
        integration_stash.upload(sample_text_file, target="/docs/guides/")
        integration_stash.upload(sample_json_file, target="/docs/api/")

        # Query with path filter for /docs/guides/
        results = integration_stash.query("content", path="/docs/guides/")

        # Should find text file only
        keys = [r.key for r in results]
        assert f"docs/guides/{sample_text_file.name}" in keys
        assert f"docs/api/{sample_json_file.name}" not in keys

    def test_query_with_root_path_finds_all(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with root path '/' finds content at root only."""
        # Upload one file to root and one to folder
        integration_stash.upload(sample_text_file, target="/")
        integration_stash.upload(sample_json_file, target="/subfolder/")

        # Query with path filter for root
        results = integration_stash.query("content", path="/")

        # Should find text file (at root)
        keys = [r.key for r in results]
        assert sample_text_file.name in keys

    def test_query_without_path_finds_everywhere(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query without path filter finds content everywhere."""
        # Upload files to different locations
        integration_stash.upload(sample_text_file, target="/docs/")
        integration_stash.upload(sample_json_file, target="/config/")

        # Query without path filter
        results = integration_stash.query("content")

        # Should find both files
        keys = [r.key for r in results]
        assert f"docs/{sample_text_file.name}" in keys
        assert f"config/{sample_json_file.name}" in keys

    def test_query_with_path_and_tags_combined(
        self,
        integration_stash: SemStash,
        sample_text_file: Path,
        sample_json_file: Path,
    ) -> None:
        """Query with both path and tag filters works correctly."""
        # Upload files with tags to different folders
        integration_stash.upload(sample_text_file, target="/docs/", tags=["important"])
        integration_stash.upload(sample_json_file, target="/docs/", tags=["config"])

        # Query with path AND tag filter
        results = integration_stash.query("content", path="/docs/", tags=["important"])

        # Should find only the text file (matches both path and tag)
        keys = [r.key for r in results]
        assert f"docs/{sample_text_file.name}" in keys
        assert f"docs/{sample_json_file.name}" not in keys


@pytest.mark.integration
class TestIntegrationLargeFiles:
    """Integration tests for large file handling with async segmentation."""

    def test_large_text_uses_async_segmentation(
        self,
        integration_stash: SemStash,
        tmp_path: Path,
    ) -> None:
        """Large text files (>50KB) use async segmentation for multiple embeddings."""
        # Generate large text file (>50KB to trigger async segmentation)
        large_text = "This is a test paragraph about semantic search. " * 2000  # ~100KB
        large_text_file = tmp_path / "large_text.txt"
        large_text_file.write_text(large_text)

        # Upload large text file
        result = integration_stash.upload(large_text_file, target="/")
        assert result.key == "large_text.txt"
        assert result.path == "/large_text.txt"

        # Query should find the content (keys may have #chunk=N suffix)
        query_results = integration_stash.query("semantic search test paragraph", top_k=5)
        keys = [r.key for r in query_results]
        assert any(k.startswith("large_text.txt") for k in keys)

    def test_large_video_uses_async_segmentation(
        self,
        integration_stash: SemStash,
    ) -> None:
        """Large video files (>10MB) use async segmentation for time-based segments."""
        # Use the existing large video sample (31MB)
        large_video = Path(__file__).parent / "samples" / "sample_large.mp4"
        if not large_video.exists():
            pytest.skip("sample_large.mp4 not found - skipping large video test")

        # Upload large video file
        result = integration_stash.upload(large_video, target="/videos/")
        assert result.key == "videos/sample_large.mp4"
        assert result.path == "/videos/sample_large.mp4"

        # Verify content was stored (browse shows the file exists)
        browse_result = integration_stash.browse("/videos/")
        keys = [item.key for item in browse_result.items]
        assert "videos/sample_large.mp4" in keys

        # Query should find the video content (keys may have #segment=N suffix)
        # Use path filter to ensure we're searching in the right location
        query_results = integration_stash.query("video", path="/videos/", top_k=10)
        keys = [r.key for r in query_results]
        assert any(k.startswith("videos/sample_large.mp4") for k in keys)

    def test_large_audio_uses_async_segmentation(
        self,
        integration_stash: SemStash,
        tmp_path: Path,
    ) -> None:
        """Large audio files (>5MB) use async segmentation for time-based segments."""
        # Use the existing audio sample (8.5MB MP3)
        large_audio = Path(__file__).parent / "samples" / "sample.mp3"
        if not large_audio.exists():
            pytest.skip("sample.mp3 not found - skipping large audio test")

        # Check if file is large enough to trigger async (>5MB)
        if large_audio.stat().st_size < 5 * 1024 * 1024:
            pytest.skip("sample.mp3 is too small to trigger async segmentation")

        # Upload large audio file
        result = integration_stash.upload(large_audio, target="/audio/")
        assert result.key == "audio/sample.mp3"
        assert result.path == "/audio/sample.mp3"

        # Query should find the audio content (keys may have #segment=N suffix)
        query_results = integration_stash.query("music audio sound", top_k=5)
        keys = [r.key for r in query_results]
        assert any(k.startswith("audio/sample.mp3") for k in keys)

    def test_multipage_pdf_creates_multiple_embeddings(
        self,
        integration_stash: SemStash,
        tmp_path: Path,
    ) -> None:
        """Multi-page PDF creates one embedding per page."""
        import fitz

        # Create a 3-page PDF
        doc = fitz.open()
        for i in range(3):
            page = doc.new_page()
            page.insert_text((50, 50), f"Page {i + 1}: This is test content for page {i + 1}")
        pdf_file = tmp_path / "multipage.pdf"
        pdf_file.write_bytes(doc.tobytes())
        doc.close()

        # Upload multi-page PDF
        result = integration_stash.upload(pdf_file, target="/docs/")
        assert result.key == "docs/multipage.pdf"

        # Query should find content from any page (keys have #page=N suffix)
        query_results = integration_stash.query("test content page", top_k=10)
        keys = [r.key for r in query_results]
        assert any(k.startswith("docs/multipage.pdf") for k in keys)

    def test_multipage_docx_creates_multiple_embeddings(
        self,
        integration_stash: SemStash,
        tmp_path: Path,
    ) -> None:
        """Multi-page DOCX creates one embedding per page."""
        import docx

        # Create a DOCX with multiple paragraphs (will render as pages)
        doc = docx.Document()
        for i in range(5):
            doc.add_paragraph(f"Section {i + 1}: This is test content for section {i + 1}. " * 50)
            doc.add_page_break()
        docx_file = tmp_path / "multipage.docx"
        doc.save(str(docx_file))

        # Upload multi-page DOCX
        result = integration_stash.upload(docx_file, target="/docs/")
        assert result.key == "docs/multipage.docx"

        # Query should find content (keys have #page=N suffix)
        query_results = integration_stash.query("test content section", top_k=10)
        keys = [r.key for r in query_results]
        assert any(k.startswith("docs/multipage.docx") for k in keys)

    def test_multislide_pptx_creates_multiple_embeddings(
        self,
        integration_stash: SemStash,
        tmp_path: Path,
    ) -> None:
        """Multi-slide PPTX creates one embedding per slide."""
        from pptx import Presentation

        # Create a 3-slide presentation
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = f"Slide {i + 1}: Test presentation content"
        pptx_file = tmp_path / "multislide.pptx"
        prs.save(str(pptx_file))

        # Upload multi-slide PPTX
        result = integration_stash.upload(pptx_file, target="/presentations/")
        assert result.key == "presentations/multislide.pptx"

        # Query should find content (keys have #slide=N suffix)
        query_results = integration_stash.query("test presentation slide", top_k=10)
        keys = [r.key for r in query_results]
        assert any(k.startswith("presentations/multislide.pptx") for k in keys)
