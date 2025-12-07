"""Nova Multimodal Embeddings generation.

Generates embeddings using Amazon Nova for text, images, audio, and video.
Supports configurable dimensions (256, 384, 1024, 3072).

Document support:
- PDF: Rendered to image using PyMuPDF, embedded with DOCUMENT_IMAGE detail level
  - Multi-page: Each page rendered and embedded separately
- DOCX: Rendered to image using PyMuPDF, embedded with DOCUMENT_IMAGE detail level
  - Multi-page: Each page rendered and embedded separately
- PPTX: Rendered to image using PyMuPDF, embedded with DOCUMENT_IMAGE detail level
  - Multi-slide: Each slide rendered and embedded separately
- XLSX: Converted to CSV text using openpyxl, embedded as text

Chunking support:
- PDF: One embedding per page (DOCUMENT_IMAGE detail level)
- DOCX: One embedding per page (DOCUMENT_IMAGE detail level)
- PPTX: One embedding per slide (DOCUMENT_IMAGE detail level)
- Large text (>50KB): Segmented via Nova Async API
- Large audio (>5MB): Segmented via Nova Async API (time-based segments)
- Large video (>10MB): Segmented via Nova Async API (time-based segments)
- Other types: Single embedding
"""

import base64
import io
import json
import mimetypes
import time
import uuid
from pathlib import Path
from typing import Any

import boto3
import docx
import fitz  # type: ignore[import-untyped]  # PyMuPDF
import openpyxl
from botocore.exceptions import ClientError
from pptx import Presentation

from semstash.config import (
    ASYNC_POLL_INTERVAL_SECONDS,
    ASYNC_POLL_MAX_INTERVAL_SECONDS,
    ASYNC_POLL_TIMEOUT_SECONDS,
    AUDIO_ASYNC_THRESHOLD_BYTES,
    DEFAULT_DIMENSION,
    DEFAULT_MEDIA_SEGMENT_SECONDS,
    DEFAULT_REGION,
    DEFAULT_TEXT_SEGMENT_CHARS,
    MAX_MEDIA_SEGMENT_SECONDS,
    MAX_TEXT_SEGMENT_CHARS,
    MIN_MEDIA_SEGMENT_SECONDS,
    MIN_TEXT_SEGMENT_CHARS,
    NOVA_EMBEDDINGS_MODEL,
    SUPPORTED_DIMENSIONS,
    TEXT_ASYNC_THRESHOLD_BYTES,
    VIDEO_ASYNC_THRESHOLD_BYTES,
)
from semstash.exceptions import (
    DimensionError,
    EmbeddingError,
    UnsupportedContentTypeError,
)
from semstash.models import (
    ChunkEmbedding,
    ChunkType,
    FileEmbeddings,
)

# Content type to modality mapping
MODALITY_MAP = {
    # Text types
    "text/plain": "text",
    "text/html": "text",
    "text/markdown": "text",
    "text/csv": "text",
    "application/json": "text",
    "application/xml": "text",
    # Image types
    "image/jpeg": "image",
    "image/png": "image",
    "image/gif": "image",
    "image/webp": "image",
    # Audio types
    "audio/mpeg": "audio",
    "audio/mp3": "audio",
    "audio/wav": "audio",
    "audio/ogg": "audio",
    "audio/flac": "audio",
    # Video types
    "video/mp4": "video",
    "video/quicktime": "video",
    "video/x-msvideo": "video",
    "video/webm": "video",
    # PDF documents (rendered to image via PyMuPDF)
    "application/pdf": "pdf",
    # Office documents (text extracted via python-docx/python-pptx)
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "office",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "office",
    # Spreadsheet types (converted to text/CSV)
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "spreadsheet",
}

# Format mappings for media types (content_type -> Nova format string)
IMAGE_FORMAT_MAP = {
    "image/jpeg": "jpeg",
    "image/jpg": "jpeg",
    "image/png": "png",
    "image/gif": "gif",
    "image/webp": "webp",
}

AUDIO_FORMAT_MAP = {
    "audio/mpeg": "mp3",
    "audio/mp3": "mp3",
    "audio/wav": "wav",
    "audio/ogg": "ogg",
    "audio/flac": "wav",  # Fallback for unsupported formats
}

VIDEO_FORMAT_MAP = {
    "video/mp4": "mp4",
    "video/quicktime": "mov",
    "video/x-msvideo": "avi",
    "video/webm": "webm",
    "video/x-matroska": "mkv",
    "video/x-flv": "flv",
    "video/mpeg": "mpeg",
    "video/3gpp": "3gp",
    "video/x-ms-wmv": "wmv",
}


class EmbeddingGenerator:
    """Generates embeddings using Amazon Nova Multimodal.

    Example:
        generator = EmbeddingGenerator(dimension=1024)

        # Generate text embedding
        embedding = generator.embed_text("Hello world")

        # Generate image embedding
        embedding = generator.embed_file(Path("photo.jpg"))

        # Get modality for a file
        modality = generator.get_modality("image/jpeg")  # "image"
    """

    def __init__(
        self,
        region: str = DEFAULT_REGION,
        dimension: int = DEFAULT_DIMENSION,
        client: Any | None = None,
    ) -> None:
        """Initialize embedding generator.

        Args:
            region: AWS region.
            dimension: Embedding dimension (256, 384, 1024, 3072).
            client: Optional boto3 bedrock-runtime client (for testing).

        Raises:
            DimensionError: If dimension is not supported.
        """
        if dimension not in SUPPORTED_DIMENSIONS:
            raise DimensionError(
                f"Invalid dimension {dimension}. "
                f"Supported: {', '.join(str(d) for d in SUPPORTED_DIMENSIONS)}"
            )

        self.region = region
        self.dimension = dimension
        self._client = client

    @property
    def client(self) -> Any:
        """Get or create Bedrock Runtime client."""
        if self._client is None:
            self._client = boto3.client("bedrock-runtime", region_name=self.region)
        return self._client

    def embed_text(self, text: str) -> list[float]:
        """Generate embedding for text content.

        Args:
            text: Text content to embed.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        request = self._build_request(
            {
                "text": {"truncationMode": "END", "value": text},
            }
        )
        return self._invoke_model(request)

    def embed_image(self, image_data: bytes, content_type: str = "image/jpeg") -> list[float]:
        """Generate embedding for image content.

        Args:
            image_data: Raw image bytes.
            content_type: MIME type of the image.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        request = self._build_request(
            {
                "image": {
                    "detailLevel": "STANDARD_IMAGE",
                    "format": IMAGE_FORMAT_MAP.get(content_type, "jpeg"),
                    "source": {"bytes": base64.b64encode(image_data).decode("utf-8")},
                },
            }
        )
        return self._invoke_model(request)

    def embed_audio(self, audio_data: bytes, content_type: str) -> list[float]:
        """Generate embedding for audio content.

        Args:
            audio_data: Raw audio bytes.
            content_type: MIME type of the audio.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        request = self._build_request(
            {
                "audio": {
                    "format": AUDIO_FORMAT_MAP.get(content_type, "mp3"),
                    "source": {"bytes": base64.b64encode(audio_data).decode("utf-8")},
                },
            }
        )
        return self._invoke_model(request)

    def embed_video(self, video_data: bytes, content_type: str) -> list[float]:
        """Generate embedding for video content.

        Args:
            video_data: Raw video bytes.
            content_type: MIME type of the video.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        request = self._build_request(
            {
                "video": {
                    "format": VIDEO_FORMAT_MAP.get(content_type, "mp4"),
                    "embeddingMode": "AUDIO_VIDEO_COMBINED",
                    "source": {"bytes": base64.b64encode(video_data).decode("utf-8")},
                },
            }
        )
        return self._invoke_model(request)

    def embed_pdf(self, pdf_data: bytes) -> list[float]:
        """Generate embedding for PDF content (first page only).

        Uses PyMuPDF to render PDF pages to images, then embeds with
        the Nova DOCUMENT_IMAGE detail level for better text interpretation.
        For multi-page documents, embeds the first page only.

        For multi-page support, use embed_pdf_pages() instead.

        Args:
            pdf_data: Raw PDF bytes.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        try:
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            if doc.page_count == 0:
                raise EmbeddingError("PDF has no pages")
            # Render first page to PNG at 2x resolution for better text recognition
            pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
            image_data = pix.tobytes("png")
            doc.close()
        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to render PDF to image: {e}") from e

        return self._embed_document_image(image_data)

    def _embed_document_image(self, image_data: bytes) -> list[float]:
        """Generate embedding for a document image (internal helper).

        Uses DOCUMENT_IMAGE detail level for better text interpretation.

        Args:
            image_data: PNG image bytes.

        Returns:
            Embedding vector as list of floats.
        """
        request = self._build_request(
            {
                "image": {
                    "detailLevel": "DOCUMENT_IMAGE",
                    "format": "png",
                    "source": {"bytes": base64.b64encode(image_data).decode("utf-8")},
                },
            }
        )
        return self._invoke_model(request)

    def embed_pdf_pages(self, pdf_data: bytes) -> FileEmbeddings:
        """Generate embeddings for all PDF pages.

        Renders each page to an image and generates a separate embedding
        using DOCUMENT_IMAGE detail level for optimal text recognition.

        Args:
            pdf_data: Raw PDF bytes.

        Returns:
            FileEmbeddings containing one ChunkEmbedding per page.

        Raises:
            EmbeddingError: If embedding generation fails.

        Example:
            result = generator.embed_pdf_pages(pdf_bytes)
            print(f"Generated {result.total_chunks} page embeddings")
            for chunk in result.chunks:
                print(f"  Page {chunk.chunk_index}: {len(chunk.embedding)} dims")
        """
        try:
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            if doc.page_count == 0:
                raise EmbeddingError("PDF has no pages")

            total_pages = doc.page_count
            chunks: list[ChunkEmbedding] = []

            for page_num in range(total_pages):
                # Render page to PNG at 2x resolution for better text recognition
                pix = doc[page_num].get_pixmap(matrix=fitz.Matrix(2, 2))
                image_data = pix.tobytes("png")

                # Generate embedding for this page
                embedding = self._embed_document_image(image_data)

                chunks.append(
                    ChunkEmbedding(
                        chunk_type=ChunkType.PAGE,
                        chunk_index=page_num + 1,  # 1-indexed
                        total_chunks=total_pages,
                        embedding=embedding,
                    )
                )

            doc.close()

            return FileEmbeddings(
                source_file="",  # Will be set by caller
                content_type="application/pdf",
                chunks=chunks,
            )

        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to embed PDF pages: {e}") from e

    def embed_office(self, office_data: bytes, content_type: str) -> list[float]:
        """Generate embedding for Office document content (DOCX, PPTX).

        Extracts text from Office documents and embeds as text.
        This captures semantic content for search without external dependencies.

        Args:
            office_data: Raw Office document bytes.
            content_type: MIME type to determine document format.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        try:
            # Content type constants for readability
            docx_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            pptx_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            extractors = {
                docx_type: self._extract_docx_text,
                pptx_type: self._extract_pptx_text,
            }
            extractor = extractors.get(content_type)
            if not extractor:
                raise EmbeddingError(f"Unsupported Office format: {content_type}")

            text = extractor(office_data)
            if not text.strip():
                raise EmbeddingError("Document contains no text")
        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to extract text from document: {e}") from e

        return self.embed_text(text)

    def _extract_docx_text(self, docx_data: bytes) -> str:
        """Extract text from DOCX document.

        Args:
            docx_data: Raw DOCX bytes.

        Returns:
            Extracted text content.
        """
        doc = docx.Document(io.BytesIO(docx_data))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs)

    def _extract_pptx_text(self, pptx_data: bytes) -> str:
        """Extract text from PPTX presentation.

        Args:
            pptx_data: Raw PPTX bytes.

        Returns:
            Extracted text content.
        """
        prs = Presentation(io.BytesIO(pptx_data))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"=== Slide {slide_num} ===")
                text_parts.extend(slide_text)

        return "\n".join(text_parts)

    def _extract_slide_text(self, slide: Any) -> str:
        """Extract text from a single PPTX slide.

        Args:
            slide: A python-pptx Slide object.

        Returns:
            Extracted text content from the slide.
        """
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
        return "\n".join(text_parts)

    def embed_pptx_slides(self, pptx_data: bytes) -> FileEmbeddings:
        """Generate embeddings for all PPTX slides by rendering each as an image.

        Uses PyMuPDF to render each slide to an image and embeds with
        DOCUMENT_IMAGE detail level for optimal text and visual recognition.

        Args:
            pptx_data: Raw PPTX bytes.

        Returns:
            FileEmbeddings containing one ChunkEmbedding per slide.

        Raises:
            EmbeddingError: If embedding generation fails.

        Example:
            result = generator.embed_pptx_slides(pptx_bytes)
            print(f"Generated {result.total_chunks} slide embeddings")
            for chunk in result.chunks:
                print(f"  Slide {chunk.chunk_index}: {len(chunk.embedding)} dims")
        """
        try:
            # PyMuPDF can open PPTX files directly
            doc = fitz.open(stream=pptx_data, filetype="pptx")
            if doc.page_count == 0:
                raise EmbeddingError("Presentation has no slides")

            total_slides = doc.page_count
            chunks: list[ChunkEmbedding] = []

            for slide_num in range(total_slides):
                # Render slide to PNG at 2x resolution for better recognition
                pix = doc[slide_num].get_pixmap(matrix=fitz.Matrix(2, 2))
                image_data = pix.tobytes("png")

                # Generate embedding using DOCUMENT_IMAGE detail level
                embedding = self._embed_document_image(image_data)

                chunks.append(
                    ChunkEmbedding(
                        chunk_type=ChunkType.SLIDE,
                        chunk_index=slide_num + 1,  # 1-indexed
                        total_chunks=total_slides,
                        embedding=embedding,
                    )
                )

            doc.close()

            return FileEmbeddings(
                source_file="",  # Will be set by caller
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                chunks=chunks,
            )

        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to embed PPTX slides: {e}") from e

    def embed_docx_pages(self, docx_data: bytes) -> FileEmbeddings:
        """Generate embeddings for DOCX document by rendering each page as an image.

        Uses PyMuPDF to render each page to an image and embeds with
        DOCUMENT_IMAGE detail level for optimal text recognition, similar to PDF handling.

        Args:
            docx_data: Raw DOCX bytes.

        Returns:
            FileEmbeddings containing one ChunkEmbedding per page.

        Raises:
            EmbeddingError: If embedding generation fails.

        Example:
            result = generator.embed_docx_pages(docx_bytes)
            print(f"Generated {result.total_chunks} page embeddings")
            for chunk in result.chunks:
                print(f"  Page {chunk.chunk_index}: {len(chunk.embedding)} dims")
        """
        try:
            # PyMuPDF can open DOCX files directly
            doc = fitz.open(stream=docx_data, filetype="docx")
            if doc.page_count == 0:
                raise EmbeddingError("Document has no pages")

            total_pages = doc.page_count
            chunks: list[ChunkEmbedding] = []

            for page_num in range(total_pages):
                # Render page to PNG at 2x resolution for better text recognition
                pix = doc[page_num].get_pixmap(matrix=fitz.Matrix(2, 2))
                image_data = pix.tobytes("png")

                # Generate embedding using DOCUMENT_IMAGE detail level
                embedding = self._embed_document_image(image_data)

                chunks.append(
                    ChunkEmbedding(
                        chunk_type=ChunkType.PAGE,
                        chunk_index=page_num + 1,  # 1-indexed
                        total_chunks=total_pages,
                        embedding=embedding,
                    )
                )

            doc.close()

            return FileEmbeddings(
                source_file="",  # Will be set by caller
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                chunks=chunks,
            )

        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to embed DOCX pages: {e}") from e

    def embed_spreadsheet(self, spreadsheet_data: bytes) -> list[float]:
        """Generate embedding for spreadsheet content (XLSX).

        Converts XLSX to CSV text format and embeds as text.
        This preserves the tabular data structure for semantic search.

        Args:
            spreadsheet_data: Raw XLSX file bytes.

        Returns:
            Embedding vector as list of floats.

        Raises:
            EmbeddingError: If embedding generation fails.
        """
        try:
            wb = openpyxl.load_workbook(io.BytesIO(spreadsheet_data), data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        text_parts.append(",".join(row_values))
            wb.close()
            text = "\n".join(text_parts)
            if not text.strip():
                raise EmbeddingError("Spreadsheet is empty")
        except EmbeddingError:
            raise
        except Exception as e:
            raise EmbeddingError(f"Failed to convert spreadsheet to text: {e}") from e

        return self.embed_text(text)

    def embed_file(self, file_path: Path, content_type: str | None = None) -> list[float]:
        """Generate embedding for a file.

        Automatically detects content type and uses appropriate embedding method.

        Args:
            file_path: Path to the file.
            content_type: Optional MIME type (auto-detected if None).

        Returns:
            Embedding vector as list of floats.

        Raises:
            UnsupportedContentTypeError: If content type is not supported.
            EmbeddingError: If embedding generation fails.
            FileNotFoundError: If file doesn't exist.
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if content_type is None:
            content_type, _ = mimetypes.guess_type(str(file_path))
            content_type = content_type or "application/octet-stream"

        modality = self.get_modality(content_type)

        # Dispatch table for modality handlers
        if modality == "text":
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return self.embed_text(text)

        # Binary modalities - read file once and dispatch
        file_data = file_path.read_bytes()

        handlers = {
            "image": lambda: self.embed_image(file_data, content_type),
            "audio": lambda: self.embed_audio(file_data, content_type),
            "video": lambda: self.embed_video(file_data, content_type),
            "pdf": lambda: self.embed_pdf(file_data),
            "office": lambda: self.embed_office(file_data, content_type),
            "spreadsheet": lambda: self.embed_spreadsheet(file_data),
        }

        handler = handlers.get(modality)
        if handler:
            return handler()

        raise UnsupportedContentTypeError(
            f"Cannot generate embedding for content type: {content_type}"
        )

    def embed_file_chunked(
        self,
        file_path: Path,
        content_type: str | None = None,
        s3_bucket: str | None = None,
        use_async: bool = True,
    ) -> FileEmbeddings:
        """Generate embeddings for a file with chunking support.

        Uses chunking for multi-page/multi-slide documents:
        - PDF: One embedding per page (DOCUMENT_IMAGE detail level)
        - PPTX: One embedding per slide (text extraction)
        - DOCX: One embedding per page (rendered to images)
        - Large text: Segmented via Nova Async API (if s3_bucket provided)
        - Large audio/video: Segmented via Nova Async API (if s3_bucket provided)
        - Other types: Single embedding wrapped in FileEmbeddings

        Args:
            file_path: Path to the file.
            content_type: Optional MIME type (auto-detected if None).
            s3_bucket: S3 bucket for async operations (required for large file segmentation).
            use_async: Whether to use async API for large files (default True).

        Returns:
            FileEmbeddings containing one or more ChunkEmbedding objects.

        Raises:
            UnsupportedContentTypeError: If content type is not supported.
            EmbeddingError: If embedding generation fails.
            FileNotFoundError: If file doesn't exist.

        Example:
            # Basic usage - auto-detect chunking strategy
            result = generator.embed_file_chunked(Path("report.pdf"))

            # With async segmentation for large files
            result = generator.embed_file_chunked(
                Path("long_document.txt"),
                s3_bucket="my-bucket",
            )

            if result.is_single_chunk:
                print("Single embedding generated")
            else:
                print(f"Generated {result.total_chunks} embeddings")
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if content_type is None:
            content_type, _ = mimetypes.guess_type(str(file_path))
            content_type = content_type or "application/octet-stream"

        modality = self.get_modality(content_type)
        file_size = file_path.stat().st_size

        # Content types that support chunking
        pptx_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        docx_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        # PDF: Embed all pages
        if modality == "pdf":
            file_data = file_path.read_bytes()
            result = self.embed_pdf_pages(file_data)
            result.source_file = file_path.name
            return result

        # PPTX: Embed all slides
        if content_type == pptx_type:
            file_data = file_path.read_bytes()
            result = self.embed_pptx_slides(file_data)
            result.source_file = file_path.name
            return result

        # DOCX: Render to pages and embed each page
        if content_type == docx_type:
            file_data = file_path.read_bytes()
            result = self.embed_docx_pages(file_data)
            result.source_file = file_path.name
            return result

        # Large text: Use async segmentation if bucket provided
        is_large_text = file_size > TEXT_ASYNC_THRESHOLD_BYTES
        if modality == "text" and use_async and s3_bucket and is_large_text:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return self.embed_text_segmented(text, s3_bucket, source_file=file_path.name)

        # Large audio: Use async segmentation if bucket provided
        is_large_audio = file_size > AUDIO_ASYNC_THRESHOLD_BYTES
        if modality == "audio" and use_async and s3_bucket and is_large_audio:
            audio_data = file_path.read_bytes()
            return self.embed_audio_segmented(
                audio_data, content_type, s3_bucket, source_file=file_path.name
            )

        # Large video: Use async segmentation if bucket provided
        is_large_video = file_size > VIDEO_ASYNC_THRESHOLD_BYTES
        if modality == "video" and use_async and s3_bucket and is_large_video:
            video_data = file_path.read_bytes()
            return self.embed_video_segmented(
                video_data, content_type, s3_bucket, source_file=file_path.name
            )

        # All other types: Single embedding wrapped in FileEmbeddings
        embedding = self.embed_file(file_path, content_type)

        return FileEmbeddings(
            source_file=file_path.name,
            content_type=content_type,
            chunks=[
                ChunkEmbedding(
                    chunk_type=ChunkType.FILE,
                    chunk_index=1,
                    total_chunks=1,
                    embedding=embedding,
                    chunk_id="",  # No fragment for single-chunk files
                )
            ],
        )

    def get_modality(self, content_type: str) -> str:
        """Determine modality from content type.

        Args:
            content_type: MIME type string.

        Returns:
            Modality string: "text", "image", "audio", "video", "pdf", "office", or "spreadsheet".

        Raises:
            UnsupportedContentTypeError: If content type is not supported.
        """
        # Check exact match
        if content_type in MODALITY_MAP:
            return MODALITY_MAP[content_type]

        # Check prefix match for text types
        if content_type.startswith("text/"):
            return "text"

        # Unknown type
        raise UnsupportedContentTypeError(
            f"Unsupported content type: {content_type}. "
            f"Supported types: {', '.join(sorted(MODALITY_MAP.keys()))}"
        )

    def _build_request(self, content_params: dict[str, Any]) -> dict[str, Any]:
        """Build Nova embedding request body.

        Args:
            content_params: Content-specific params (text, image, audio, or video).

        Returns:
            Complete request body for invoke_model.
        """
        return {
            "taskType": "SINGLE_EMBEDDING",
            "singleEmbeddingParams": {
                "embeddingPurpose": "GENERIC_INDEX",
                "embeddingDimension": self.dimension,
                **content_params,
            },
        }

    def _invoke_model(self, request_body: dict[str, Any]) -> list[float]:
        """Invoke the Nova embeddings model.

        Args:
            request_body: Model request payload.

        Returns:
            Embedding vector.

        Raises:
            EmbeddingError: If model invocation fails.
        """
        try:
            response = self.client.invoke_model(
                modelId=NOVA_EMBEDDINGS_MODEL,
                body=json.dumps(request_body),
                contentType="application/json",
                accept="application/json",
            )

            response_body = json.loads(response["body"].read())

            # New response format: {"embeddings": [{"embeddingType": "...", "embedding": [...]}]}
            embeddings = response_body.get("embeddings")
            if not embeddings:
                raise EmbeddingError("No embeddings in model response")

            embedding: list[float] | None = embeddings[0].get("embedding")
            if embedding is None:
                raise EmbeddingError("No embedding vector in model response")

            return embedding

        except ClientError as e:
            error_code = e.response.get("Error", {}).get("Code", "Unknown")
            error_msg = e.response.get("Error", {}).get("Message", str(e))
            raise EmbeddingError(f"Bedrock error ({error_code}): {error_msg}") from e

        except json.JSONDecodeError as e:
            raise EmbeddingError(f"Invalid model response: {e}") from e

    # -------------------------------------------------------------------------
    # Nova Async API - Segmented Embeddings
    # -------------------------------------------------------------------------

    def _get_s3_client(self) -> Any:
        """Get or create S3 client for async operations."""
        if not hasattr(self, "_s3_client"):
            self._s3_client = boto3.client("s3", region_name=self.region)
        return self._s3_client

    def _start_async_invoke(
        self,
        input_s3_uri: str,
        output_s3_uri: str,
        modality: str,
        segment_config: dict[str, Any],
        media_format: str | None = None,
    ) -> str:
        """Start an async embedding job via Nova API.

        Args:
            input_s3_uri: S3 URI for input content.
            output_s3_uri: S3 URI prefix for output.
            modality: Content modality ("text", "audio", "video").
            segment_config: Segmentation configuration.
            media_format: Format for audio/video (e.g., "mp3", "mp4").

        Returns:
            Job ID from Nova.

        Raises:
            EmbeddingError: If job submission fails.
        """
        # Build modality-specific params with source and segmentation config
        modality_params: dict[str, Any] = {
            "source": {"s3Location": {"uri": input_s3_uri}},
            "segmentationConfig": segment_config,
        }

        # Text requires truncationMode
        if modality == "text":
            modality_params["truncationMode"] = "END"

        # Video requires embeddingMode and format
        if modality == "video":
            modality_params["embeddingMode"] = "AUDIO_VIDEO_COMBINED"
            if media_format:
                modality_params["format"] = media_format

        # Audio requires format
        if modality == "audio" and media_format:
            modality_params["format"] = media_format

        request_body: dict[str, Any] = {
            "taskType": "SEGMENTED_EMBEDDING",
            "segmentedEmbeddingParams": {
                "embeddingPurpose": "GENERIC_INDEX",
                "embeddingDimension": self.dimension,
                modality: modality_params,
            },
        }

        try:
            response = self.client.start_async_invoke(
                modelId=NOVA_EMBEDDINGS_MODEL,
                modelInput=request_body,
                outputDataConfig={"s3OutputDataConfig": {"s3Uri": output_s3_uri}},
            )
            return str(response["invocationArn"])
        except ClientError as e:
            error_code = e.response.get("Error", {}).get("Code", "Unknown")
            error_msg = e.response.get("Error", {}).get("Message", str(e))
            raise EmbeddingError(f"Failed to start async job ({error_code}): {error_msg}") from e

    def _poll_async_invoke(self, invocation_arn: str) -> tuple[str, str]:
        """Poll for async job completion.

        Args:
            invocation_arn: The invocation ARN from start_async_invoke.

        Returns:
            Tuple of (status, output_s3_uri or error_message).

        Raises:
            EmbeddingError: If polling fails.
        """
        try:
            response = self.client.get_async_invoke(invocationArn=invocation_arn)
            status = response.get("status", "Unknown")

            if status == "Completed":
                output_config = response.get("outputDataConfig", {})
                s3_config = output_config.get("s3OutputDataConfig", {})
                output_uri = s3_config.get("s3Uri", "")
                return ("Completed", output_uri)
            elif status == "Failed":
                failure_msg = response.get("failureMessage", "Unknown error")
                return ("Failed", failure_msg)
            else:
                return (status, "")
        except ClientError as e:
            error_code = e.response.get("Error", {}).get("Code", "Unknown")
            error_msg = e.response.get("Error", {}).get("Message", str(e))
            raise EmbeddingError(f"Failed to poll async job ({error_code}): {error_msg}") from e

    def _parse_s3_uri(self, s3_uri: str) -> tuple[str, str]:
        """Parse S3 URI into bucket and key.

        Args:
            s3_uri: S3 URI like 's3://bucket/key'.

        Returns:
            Tuple of (bucket, key).
        """
        if not s3_uri.startswith("s3://"):
            raise ValueError(f"Invalid S3 URI: {s3_uri}")
        path = s3_uri[5:]  # Remove 's3://'
        bucket, _, key = path.partition("/")
        return (bucket, key)

    def _read_async_output(self, output_s3_uri: str) -> list[dict[str, Any]]:
        """Read and parse JSONL output from async job.

        Args:
            output_s3_uri: S3 URI for output directory.

        Returns:
            List of embedding result dictionaries.

        Raises:
            EmbeddingError: If reading or parsing fails.
        """
        try:
            s3 = self._get_s3_client()
            bucket, prefix = self._parse_s3_uri(output_s3_uri)

            # List objects in output prefix to find the JSONL file
            response = s3.list_objects_v2(Bucket=bucket, Prefix=prefix)
            contents = response.get("Contents", [])

            if not contents:
                raise EmbeddingError(f"No output files found at {output_s3_uri}")

            # Find the embedding JSONL file (e.g., embedding-text.jsonl)
            jsonl_key = None
            for obj in contents:
                key = obj["Key"]
                if key.endswith(".jsonl") and "embedding-" in key:
                    jsonl_key = key
                    break

            if not jsonl_key:
                # Fallback: try any .jsonl file
                for obj in contents:
                    if obj["Key"].endswith(".jsonl"):
                        jsonl_key = obj["Key"]
                        break

            if not jsonl_key:
                raise EmbeddingError(f"No JSONL output file found at {output_s3_uri}")

            # Read and parse the file
            obj_response = s3.get_object(Bucket=bucket, Key=jsonl_key)
            body = obj_response["Body"].read().decode("utf-8")

            results = []
            for line in body.strip().split("\n"):
                if line.strip():
                    results.append(json.loads(line))

            return results

        except ClientError as e:
            error_code = e.response.get("Error", {}).get("Code", "Unknown")
            error_msg = e.response.get("Error", {}).get("Message", str(e))
            raise EmbeddingError(f"Failed to read async output ({error_code}): {error_msg}") from e
        except json.JSONDecodeError as e:
            raise EmbeddingError(f"Failed to parse async output: {e}") from e

    def embed_text_segmented(
        self,
        text: str,
        s3_bucket: str,
        s3_prefix: str = ".semstash/async/",
        segment_chars: int = DEFAULT_TEXT_SEGMENT_CHARS,
        source_file: str = "",
    ) -> FileEmbeddings:
        """Generate segmented embeddings for text using Nova Async API.

        Uses Nova's native text segmentation to create multiple embeddings
        for long text content. Each segment is embedded separately, providing
        better semantic coverage than truncation.

        Args:
            text: Text content to embed.
            s3_bucket: S3 bucket for temporary input/output.
            s3_prefix: Prefix for temporary files (default: .semstash/async/).
            segment_chars: Max characters per segment (300-50000).
            source_file: Original filename for result metadata.

        Returns:
            FileEmbeddings with one ChunkEmbedding per segment.

        Raises:
            EmbeddingError: If embedding generation fails.
            ValueError: If segment_chars is out of range.

        Example:
            result = generator.embed_text_segmented(
                long_text,
                s3_bucket="my-bucket",
                segment_chars=5000,
            )
            print(f"Generated {result.total_chunks} segment embeddings")
        """
        if segment_chars < MIN_TEXT_SEGMENT_CHARS or segment_chars > MAX_TEXT_SEGMENT_CHARS:
            raise ValueError(
                f"segment_chars must be {MIN_TEXT_SEGMENT_CHARS}-{MAX_TEXT_SEGMENT_CHARS}, "
                f"got {segment_chars}"
            )

        # Generate unique job ID
        job_id = str(uuid.uuid4())[:8]
        input_key = f"{s3_prefix.rstrip('/')}/input/{job_id}.txt"
        output_prefix = f"{s3_prefix.rstrip('/')}/output/{job_id}/"

        try:
            # Upload text to S3
            s3 = self._get_s3_client()
            s3.put_object(Bucket=s3_bucket, Key=input_key, Body=text.encode("utf-8"))

            input_s3_uri = f"s3://{s3_bucket}/{input_key}"
            output_s3_uri = f"s3://{s3_bucket}/{output_prefix}"

            # Start async job
            segment_config = {"maxLengthChars": segment_chars}
            invocation_arn = self._start_async_invoke(
                input_s3_uri, output_s3_uri, "text", segment_config
            )

            # Poll for completion with exponential backoff
            interval = ASYNC_POLL_INTERVAL_SECONDS
            elapsed = 0.0

            while elapsed < ASYNC_POLL_TIMEOUT_SECONDS:
                time.sleep(interval)
                elapsed += interval

                status, result = self._poll_async_invoke(invocation_arn)

                if status == "Completed":
                    # Parse results
                    raw_results = self._read_async_output(result)
                    return self._parse_segmented_results(
                        raw_results, ChunkType.CHUNK, "text/plain", source_file
                    )
                elif status == "Failed":
                    raise EmbeddingError(f"Async job failed: {result}")

                # Exponential backoff
                interval = min(interval * 1.5, ASYNC_POLL_MAX_INTERVAL_SECONDS)

            raise EmbeddingError(f"Async job timed out after {ASYNC_POLL_TIMEOUT_SECONDS}s")

        finally:
            # Cleanup temporary files
            try:
                s3 = self._get_s3_client()
                s3.delete_object(Bucket=s3_bucket, Key=input_key)
                # List and delete output files
                response = s3.list_objects_v2(Bucket=s3_bucket, Prefix=output_prefix)
                for obj in response.get("Contents", []):
                    s3.delete_object(Bucket=s3_bucket, Key=obj["Key"])
            except Exception:
                pass  # Best effort cleanup

    def embed_audio_segmented(
        self,
        audio_data: bytes,
        content_type: str,
        s3_bucket: str,
        s3_prefix: str = ".semstash/async/",
        segment_seconds: int = DEFAULT_MEDIA_SEGMENT_SECONDS,
        source_file: str = "",
    ) -> FileEmbeddings:
        """Generate segmented embeddings for audio using Nova Async API.

        Uses Nova's native audio segmentation to create multiple embeddings
        for long audio content. Each segment covers a time window.

        Args:
            audio_data: Raw audio bytes.
            content_type: MIME type of audio.
            s3_bucket: S3 bucket for temporary input/output.
            s3_prefix: Prefix for temporary files.
            segment_seconds: Duration per segment in seconds (5-30).
            source_file: Original filename for result metadata.

        Returns:
            FileEmbeddings with one ChunkEmbedding per time segment.

        Raises:
            EmbeddingError: If embedding generation fails.
            ValueError: If segment_seconds is below minimum.

        Example:
            result = generator.embed_audio_segmented(
                audio_bytes,
                "audio/mp3",
                s3_bucket="my-bucket",
                segment_seconds=30,
            )
            print(f"Generated {result.total_chunks} audio segment embeddings")
        """
        if not MIN_MEDIA_SEGMENT_SECONDS <= segment_seconds <= MAX_MEDIA_SEGMENT_SECONDS:
            raise ValueError(
                f"segment_seconds must be {MIN_MEDIA_SEGMENT_SECONDS}-{MAX_MEDIA_SEGMENT_SECONDS}, "
                f"got {segment_seconds}"
            )

        audio_format = AUDIO_FORMAT_MAP.get(content_type, "mp3")
        job_id = str(uuid.uuid4())[:8]
        input_key = f"{s3_prefix.rstrip('/')}/input/{job_id}.{audio_format}"
        output_prefix = f"{s3_prefix.rstrip('/')}/output/{job_id}/"

        try:
            # Upload audio to S3
            s3 = self._get_s3_client()
            s3.put_object(Bucket=s3_bucket, Key=input_key, Body=audio_data)

            input_s3_uri = f"s3://{s3_bucket}/{input_key}"
            output_s3_uri = f"s3://{s3_bucket}/{output_prefix}"

            # Start async job with audio-specific config
            segment_config = {"durationSeconds": segment_seconds}
            invocation_arn = self._start_async_invoke(
                input_s3_uri, output_s3_uri, "audio", segment_config, media_format=audio_format
            )

            # Poll for completion
            interval = ASYNC_POLL_INTERVAL_SECONDS
            elapsed = 0.0

            while elapsed < ASYNC_POLL_TIMEOUT_SECONDS:
                time.sleep(interval)
                elapsed += interval

                status, result = self._poll_async_invoke(invocation_arn)

                if status == "Completed":
                    raw_results = self._read_async_output(result)
                    return self._parse_segmented_results(
                        raw_results, ChunkType.SEGMENT, content_type, source_file
                    )
                elif status == "Failed":
                    raise EmbeddingError(f"Async job failed: {result}")

                interval = min(interval * 1.5, ASYNC_POLL_MAX_INTERVAL_SECONDS)

            raise EmbeddingError(f"Async job timed out after {ASYNC_POLL_TIMEOUT_SECONDS}s")

        finally:
            # Cleanup
            try:
                s3 = self._get_s3_client()
                s3.delete_object(Bucket=s3_bucket, Key=input_key)
                response = s3.list_objects_v2(Bucket=s3_bucket, Prefix=output_prefix)
                for obj in response.get("Contents", []):
                    s3.delete_object(Bucket=s3_bucket, Key=obj["Key"])
            except Exception:
                pass

    def embed_video_segmented(
        self,
        video_data: bytes,
        content_type: str,
        s3_bucket: str,
        s3_prefix: str = ".semstash/async/",
        segment_seconds: int = DEFAULT_MEDIA_SEGMENT_SECONDS,
        source_file: str = "",
    ) -> FileEmbeddings:
        """Generate segmented embeddings for video using Nova Async API.

        Uses Nova's native video segmentation to create multiple embeddings
        for long video content. Each segment covers a time window and includes
        both audio and visual information.

        Args:
            video_data: Raw video bytes.
            content_type: MIME type of video.
            s3_bucket: S3 bucket for temporary input/output.
            s3_prefix: Prefix for temporary files.
            segment_seconds: Duration per segment in seconds (5-30).
            source_file: Original filename for result metadata.

        Returns:
            FileEmbeddings with one ChunkEmbedding per time segment.

        Raises:
            EmbeddingError: If embedding generation fails.
            ValueError: If segment_seconds is below minimum.

        Example:
            result = generator.embed_video_segmented(
                video_bytes,
                "video/mp4",
                s3_bucket="my-bucket",
                segment_seconds=30,
            )
            print(f"Generated {result.total_chunks} video segment embeddings")
        """
        if not MIN_MEDIA_SEGMENT_SECONDS <= segment_seconds <= MAX_MEDIA_SEGMENT_SECONDS:
            raise ValueError(
                f"segment_seconds must be {MIN_MEDIA_SEGMENT_SECONDS}-{MAX_MEDIA_SEGMENT_SECONDS}, "
                f"got {segment_seconds}"
            )

        video_format = VIDEO_FORMAT_MAP.get(content_type, "mp4")
        job_id = str(uuid.uuid4())[:8]
        input_key = f"{s3_prefix.rstrip('/')}/input/{job_id}.{video_format}"
        output_prefix = f"{s3_prefix.rstrip('/')}/output/{job_id}/"

        try:
            # Upload video to S3
            s3 = self._get_s3_client()
            s3.put_object(Bucket=s3_bucket, Key=input_key, Body=video_data)

            input_s3_uri = f"s3://{s3_bucket}/{input_key}"
            output_s3_uri = f"s3://{s3_bucket}/{output_prefix}"

            # Start async job with video-specific config
            segment_config = {"durationSeconds": segment_seconds}
            invocation_arn = self._start_async_invoke(
                input_s3_uri, output_s3_uri, "video", segment_config, media_format=video_format
            )

            # Poll for completion
            interval = ASYNC_POLL_INTERVAL_SECONDS
            elapsed = 0.0

            while elapsed < ASYNC_POLL_TIMEOUT_SECONDS:
                time.sleep(interval)
                elapsed += interval

                status, result = self._poll_async_invoke(invocation_arn)

                if status == "Completed":
                    raw_results = self._read_async_output(result)
                    return self._parse_segmented_results(
                        raw_results, ChunkType.SEGMENT, content_type, source_file
                    )
                elif status == "Failed":
                    raise EmbeddingError(f"Async job failed: {result}")

                interval = min(interval * 1.5, ASYNC_POLL_MAX_INTERVAL_SECONDS)

            raise EmbeddingError(f"Async job timed out after {ASYNC_POLL_TIMEOUT_SECONDS}s")

        finally:
            # Cleanup
            try:
                s3 = self._get_s3_client()
                s3.delete_object(Bucket=s3_bucket, Key=input_key)
                response = s3.list_objects_v2(Bucket=s3_bucket, Prefix=output_prefix)
                for obj in response.get("Contents", []):
                    s3.delete_object(Bucket=s3_bucket, Key=obj["Key"])
            except Exception:
                pass

    def _parse_segmented_results(
        self,
        raw_results: list[dict[str, Any]],
        chunk_type: ChunkType,
        content_type: str,
        source_file: str,
    ) -> FileEmbeddings:
        """Parse raw segmented embedding results into FileEmbeddings.

        Args:
            raw_results: List of result dicts from JSONL output.
            chunk_type: Type of chunks (CHUNK for text, SEGMENT for audio/video).
            content_type: MIME type of source content.
            source_file: Original filename.

        Returns:
            FileEmbeddings with parsed ChunkEmbedding objects.
        """
        if not raw_results:
            raise EmbeddingError("No embedding results returned")

        total_chunks = len(raw_results)
        chunks: list[ChunkEmbedding] = []

        for i, result in enumerate(raw_results, 1):
            # Check status
            status = result.get("status", "SUCCESS")
            if status == "FAILURE":
                reason = result.get("failureReason", "Unknown")
                msg = result.get("message", "")
                raise EmbeddingError(f"Segment {i} failed: {reason} - {msg}")

            # New format: {"embedding": [...], "segmentMetadata": {...}}
            embedding = result.get("embedding")
            if not embedding:
                # Fallback: old format {"embeddings": [{"embedding": [...]}]}
                embeddings_data = result.get("embeddings", [])
                if embeddings_data:
                    embedding = embeddings_data[0].get("embedding")

            if not embedding:
                raise EmbeddingError(f"No embedding vector in result {i}")

            chunks.append(
                ChunkEmbedding(
                    chunk_type=chunk_type,
                    chunk_index=i,
                    total_chunks=total_chunks,
                    embedding=embedding,
                )
            )

        return FileEmbeddings(
            source_file=source_file,
            content_type=content_type,
            chunks=chunks,
        )


def get_supported_content_types() -> list[str]:
    """Get list of all supported content types for embedding.

    Returns:
        Sorted list of MIME type strings.
    """
    return sorted(MODALITY_MAP.keys())


def is_supported_content_type(content_type: str) -> bool:
    """Check if a content type is supported for embedding.

    Args:
        content_type: MIME type string.

    Returns:
        True if supported, False otherwise.
    """
    if content_type in MODALITY_MAP:
        return True
    return content_type.startswith("text/")
